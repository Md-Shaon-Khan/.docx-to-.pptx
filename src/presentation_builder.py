# src/presentation_builder.py
import os
import random
import re
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

from src.ai_processor import AIProcessor
from src.image_fetcher import ImageFetcher

logger = logging.getLogger(__name__)

class PresentationBuilder:
    def __init__(self, template_path=None):
        """Initialize presentation, optionally with a template"""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.ai = AIProcessor()
        self.image_fetcher = ImageFetcher()
        self.default_layout = self.prs.slide_layouts[1]  # Title + Content layout
        
        # Define slide layouts
        self.layouts = {
            'title': 0,        # Title slide
            'title_content': 1, # Title and content
            'section_header': 2, # Section header
            'two_content': 3,  # Two content
            'comparison': 4,   # Comparison
            'blank': 5,        # Blank
            'content_caption': 6, # Content with caption
            'picture_caption': 7  # Picture with caption
        }

    def add_slide(self, title="", content="", image_query=None, layout_type='title_content', slide_number=None):
        """Add a slide with title, content, and optional image"""
        try:
            layout_index = self.layouts.get(layout_type, 1)
            slide_layout = self.prs.slide_layouts[layout_index]
            slide = self.prs.slides.add_slide(slide_layout)

            # Set title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
                for paragraph in title_shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(44 if layout_type == 'title' else 36)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

            # Set content
            content_placeholder = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_placeholder = shape
                    break
            
            if content_placeholder and content:
                # Clean and format content
                content = content.strip()
                if not content.endswith('.'):
                    content += '.'
                
                content_placeholder.text = content
                for paragraph in content_placeholder.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.line_spacing = 1.2
                    for run in paragraph.runs:
                        run.font.size = Pt(24 if layout_type == 'title' else 20)
                        run.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

            # Add colorful background based on slide type
            self._add_slide_background(slide, layout_type)

            # Add slide number if provided
            if slide_number is not None:
                self._add_slide_number(slide, slide_number)

            # Add image if provided
            if image_query:
                self._add_image_to_slide(slide, image_query, title)

            # Add decorative elements for better design
            self._add_decorative_elements(slide, layout_type)

            return slide
        
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
            return None

    def _add_slide_background(self, slide, layout_type):
        """Add colorful background to slide"""
        try:
            if layout_type == 'title':
                colors = [(41, 128, 185), (52, 152, 219), (41, 128, 185)]  # Blue shades
            elif layout_type == 'section_header':
                colors = [(39, 174, 96), (46, 204, 113), (39, 174, 96)]  # Green shades
            else:
                colors = [(255, 179, 186), (255, 223, 186), (255, 255, 186),
                         (186, 255, 201), (186, 225, 255), (210, 180, 255)]
            
            r, g, b = random.choice(colors)
            
            # Create gradient-like effect
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
            
            # Add subtle pattern if not title slide
            if layout_type != 'title':
                self._add_background_pattern(slide, r, g, b)
                
        except Exception as e:
            logger.error(f"Error adding background: {e}")

    def _add_background_pattern(self, slide, r, g, b):
        """Add subtle pattern to background"""
        try:
            # Create a text box with subtle dots pattern
            from pptx.util import Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            # Add subtle dots in corners
            for i in range(20):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    random.randint(0, 1200),
                    random.randint(0, 800),
                    Pt(10), Pt(10)
                )
                shape.fill.solid()
                # Lighter or darker shade
                if random.choice([True, False]):
                    shape.fill.fore_color.rgb = RGBColor(
                        min(255, r + 40),
                        min(255, g + 40),
                        min(255, b + 40)
                    )
                else:
                    shape.fill.fore_color.rgb = RGBColor(
                        max(0, r - 40),
                        max(0, g - 40),
                        max(0, b - 40)
                    )
                shape.line.fill.background()  # No border
        except:
            pass  # Pattern is optional

    def _add_slide_number(self, slide, number):
        """Add slide number to bottom right corner"""
        try:
            from pptx.util import Inches, Pt
            
            left = Inches(9.5)
            top = Inches(6.8)
            width = Inches(0.5)
            height = Inches(0.3)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(number)
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(100, 100, 100)
            p.alignment = PP_ALIGN.RIGHT
        except Exception as e:
            logger.error(f"Error adding slide number: {e}")

    def _add_image_to_slide(self, slide, image_query, title):
        """Add image to slide based on query"""
        try:
            os.makedirs("uploads/temp", exist_ok=True)
            # sanitize title to be filesystem-safe
            safe_title = re.sub(r'[^A-Za-z0-9._-]', '_', (title or '')).strip('_')
            if not safe_title:
                safe_title = 'image'
            image_path = f"uploads/temp/{safe_title[:50]}_{random.randint(1000, 9999)}.jpg"
            
            # Generate better image query using AI
            ai_query = self.ai.generate_image_query(image_query)
            
            result = self.image_fetcher.fetch_image(ai_query, image_path)
            if result and os.path.exists(image_path):
                # Randomize image position for variety
                positions = [
                    (Inches(5.5), Inches(1.5), Inches(4), Inches(3)),  # Right side
                    (Inches(1), Inches(1.5), Inches(4), Inches(3)),    # Left side
                    (Inches(3), Inches(4), Inches(4), Inches(2.5)),    # Bottom center
                ]
                left, top, width, height = random.choice(positions)
                
                # Add rounded corners effect by adding a white border
                pic = slide.shapes.add_picture(image_path, left, top, width, height)
                
                # Add a subtle shadow effect
                # set shadow where supported; avoid setting color attribute (not available in some pptx versions)
                try:
                    pic.shadow.inherit = False
                    pic.shadow.visible = True
                    pic.shadow.blur_radius = Pt(5)
                    pic.shadow.distance = Pt(3)
                    pic.shadow.angle = 45
                except Exception:
                    pass
                
        except Exception as e:
            logger.error(f"Error adding image: {e}")

    def _add_decorative_elements(self, slide, layout_type):
        """Add decorative elements to slide"""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            
            # Add decorative shapes based on layout
            if layout_type == 'title':
                # Add a decorative line under title
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1), Inches(2.5),
                    Inches(8), Pt(4)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold
                line.line.fill.background()
                
            elif layout_type == 'section_header':
                # Add corner decorations
                for corner in [(0, 0), (9.5, 0), (0, 6.8), (9.5, 6.8)]:
                    left, top = corner
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(left), Inches(top),
                        Inches(0.5), Inches(0.5)
                    )
                    shape.fill.solid()
                    # python-pptx RGBColor accepts 3 args; do not pass alpha
                    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
                    shape.line.fill.background()
                    
        except Exception as e:
            logger.error(f"Error adding decorative elements: {e}")

    def add_slides_from_document(self, headings, sections):
        """
        Create slides using headings and per-heading sections.
        `headings` is a list of dicts {'text', 'level'} and `sections`
        is a list of lists where each inner list contains paragraphs
        for the corresponding heading. This enables multiple slides
        per heading via AI splitting.
        """
        try:
            slide_number = 1
            
            # 1. Title Slide
            self.add_slide(
                title="Presentation",
                content="Generated from Word Document\nUsing AI-Powered Conversion",
                layout_type='title',
                slide_number=slide_number
            )
            slide_number += 1

            # If no headings found, create slides from paragraphs directly
            if not headings:
                # sections may be a flat list in older callers; accept both
                if isinstance(sections, list) and sections and isinstance(sections[0], list):
                    # no headings but sections exist
                    paras = [p for sec in sections for p in sec]
                else:
                    paras = sections if isinstance(sections, list) else []
                self._create_slides_from_paragraphs_only(paras, slide_number)
                return

            # 2. Table of Contents Slide (if we have multiple headings)
            if len(headings) > 1:
                toc_content = "Agenda:\n\n"
                for i, heading in enumerate(headings[:5]):  # Limit to 5 items
                    toc_content += f"• {heading['text']}\n"
                if len(headings) > 5:
                    toc_content += f"• ... and {len(headings) - 5} more sections\n"
                
                self.add_slide(
                    title="Agenda",
                    content=toc_content,
                    layout_type='section_header',
                    slide_number=slide_number,
                    image_query="agenda planning"
                )
                slide_number += 1

            # 3. Process content by sections
            # sections is expected to be a list of lists aligned with headings
            for heading_idx, heading in enumerate(headings):
                heading_text = heading['text']

                if heading_idx < len(sections):
                    section_content = sections[heading_idx]
                else:
                    section_content = []
                
                # If no content under this heading, create a simple slide
                if not section_content:
                    self.add_slide(
                        title=heading_text,
                        content=f"Discussion about {heading_text.lower()}",
                        layout_type='title_content',
                        slide_number=slide_number,
                        image_query=heading_text
                    )
                    slide_number += 1
                    continue
                
                # Combine all content for this section
                full_section_text = "\n".join(section_content)
                
                # Use AI to analyze and split content
                slide_contents = self._ai_split_content(full_section_text, heading_text)
                
                # Create slides for this section
                for i, (slide_title, slide_content) in enumerate(slide_contents):
                    # Use section header for first slide of section
                    layout = 'section_header' if i == 0 else 'title_content'
                    
                    # Add slide number indicator for multi-slide sections
                    slide_title_formatted = slide_title
                    if len(slide_contents) > 1:
                        slide_title_formatted = f"{heading_text} ({i+1}/{len(slide_contents)})"
                    
                    self.add_slide(
                        title=slide_title_formatted,
                        content=slide_content,
                        layout_type=layout,
                        slide_number=slide_number,
                        image_query=heading_text if i == 0 else f"{heading_text} part {i+1}"
                    )
                    slide_number += 1
            
            # 4. Summary/Thank You Slide
            self.add_slide(
                title="Summary & Next Steps",
                content=self._generate_summary_content(headings),
                layout_type='title_content',
                slide_number=slide_number,
                image_query="conclusion summary"
            )
            slide_number += 1
            
            # 5. Final Thank You Slide
            self.add_slide(
                title="Thank You",
                content="Questions & Discussion",
                layout_type='title',
                slide_number=slide_number
            )
            
        except Exception as e:
            logger.error(f"Error creating slides from document: {e}")
            import traceback
            traceback.print_exc()

    def _create_slides_from_paragraphs_only(self, paragraphs, start_slide_number):
        """Create slides when document has no headings"""
        slide_number = start_slide_number
        
        # Combine all paragraphs
        full_text = "\n".join(paragraphs)
        
        # Split into chunks for slides
        max_chars_per_slide = 800
        chunks = []
        current_chunk = ""
        
        for para in paragraphs:
            if len(current_chunk) + len(para) + 1 <= max_chars_per_slide:
                current_chunk += para + "\n"
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = para + "\n"
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Create slides for each chunk
        for i, chunk in enumerate(chunks):
            # Use AI to generate a title and summarize
            slide_data = self._ai_process_chunk(chunk, i, len(chunks))
            
            self.add_slide(
                title=slide_data['title'],
                content=slide_data['content'],
                layout_type='title_content',
                slide_number=slide_number,
                image_query=slide_data['image_query']
            )
            slide_number += 1

    def _ai_split_content(self, content, section_title):
        """
        Use AI to split content into logical slides
        Returns list of (slide_title, slide_content) tuples
        """
        try:
            prompt = f"""Split this content into logical presentation slides for the section "{section_title}".
            
            Content to split:
            {content[:2000]}  # Limit content length
            
            Provide 2-4 slide outlines with:
            1. A short slide title (not the section title)
            2. Concise content for each slide (2-4 bullet points maximum per slide)
            
            Format as:
            Slide 1 Title: [title]
            Content: [2-4 bullet points]
            
            Slide 2 Title: [title]
            Content: [2-4 bullet points]
            
            ..."""
            
            # Use AIProcessor wrapper to create chat completions
            response = self.ai.create_chat_completion(
                messages=[
                    {"role": "system", "content": "You are a presentation design expert. Split content into logical slides with clear titles and concise bullet points."},
                    {"role": "user", "content": prompt}
                ],
                model="gpt-3.5-turbo",
                max_tokens=500,
                temperature=0.7
            )

            try:
                ai_response = response.choices[0].message.content
            except Exception:
                ai_response = response['choices'][0]['message']['content']
            
            # Parse the AI response
            slides = []
            lines = ai_response.strip().split('\n')
            current_title = None
            current_content = []
            
            for line in lines:
                line = line.strip()
                if line.startswith('Slide') and 'Title:' in line:
                    # Save previous slide
                    if current_title and current_content:
                        slides.append((current_title, '\n'.join(current_content)))
                    
                    # Start new slide
                    current_title = line.split('Title:')[1].strip()
                    current_content = []
                elif line.startswith('Content:') or line.startswith('•') or line.startswith('-'):
                    if line.startswith('Content:'):
                        line = line[8:].strip()
                    if line:
                        current_content.append(line)
                elif line and not line.startswith('---') and not line.startswith('=='):
                    # Add as content line
                    current_content.append(line)
            
            # Add the last slide
            if current_title and current_content:
                slides.append((current_title, '\n'.join(current_content)))
            
            # If AI couldn't split properly, create default splits
            if not slides:
                slides = self._default_content_split(content, section_title)
            
            return slides
            
        except Exception as e:
            logger.error(f"Error in AI content splitting: {e}")
            return self._default_content_split(content, section_title)

    def _default_content_split(self, content, section_title):
        """Default method to split content when AI fails"""
        slides = []
        
        # Split by sentences
        sentences = content.replace('\n', ' ').split('. ')
        
        # Group sentences into slides
        max_sentences_per_slide = 4
        for i in range(0, len(sentences), max_sentences_per_slide):
            slide_sentences = sentences[i:i+max_sentences_per_slide]
            slide_content = '• ' + '\n• '.join(slide_sentences)
            
            # Generate slide title
            if i == 0:
                slide_title = section_title
            else:
                slide_title = f"{section_title} - Part {i//max_sentences_per_slide + 1}"
            
            slides.append((slide_title, slide_content))
        
        return slides

    def _ai_process_chunk(self, chunk, chunk_index, total_chunks):
        """Use AI to process a text chunk into slide content"""
        try:
            prompt = f"""Create presentation slide content from this text (Part {chunk_index+1}/{total_chunks}):
            
            {chunk[:1000]}
            
            Provide:
            1. A catchy slide title (3-7 words)
            2. Concise bullet points (3-5 points)
            3. A relevant image search query (2-3 words)
            
            Format as:
            Title: [title]
            Content: [bullet points]
            Image Query: [query]"""
            
            response = self.ai.create_chat_completion(
                messages=[
                    {"role": "system", "content": "You create engaging presentation slide content from text."},
                    {"role": "user", "content": prompt}
                ],
                model="gpt-3.5-turbo",
                max_tokens=300,
                temperature=0.7
            )

            try:
                result = response.choices[0].message.content
            except Exception:
                result = response['choices'][0]['message']['content']
            
            # Parse the response
            title = "Key Points"
            content = "• " + chunk[:100].replace('\n', '\n• ')
            image_query = "presentation"
            
            for line in result.split('\n'):
                if line.startswith('Title:'):
                    title = line.split('Title:')[1].strip()
                elif line.startswith('Content:'):
                    content = line.split('Content:')[1].strip()
                elif line.startswith('Image Query:'):
                    image_query = line.split('Image Query:')[1].strip()
            
            return {
                'title': title,
                'content': content,
                'image_query': image_query
            }
            
        except Exception as e:
            logger.error(f"Error processing chunk with AI: {e}")
            return {
                'title': f"Section {chunk_index + 1}",
                'content': "• " + chunk[:200].replace('\n', '\n• '),
                'image_query': "information"
            }

    def _generate_summary_content(self, headings):
        """Generate summary content based on headings"""
        if not headings:
            return "• Review of key points\n• Next steps\n• Questions & Answers"
        
        summary = "Key Takeaways:\n\n"
        for i, heading in enumerate(headings[:4]):  # Limit to 4 main points
            summary += f"• {heading['text']}\n"
        
        summary += "\nNext Steps:\n• Implement learnings\n• Schedule follow-up\n• Share feedback"
        
        return summary

    def save(self, output_path: str):
        """Save the presentation to disk"""
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.prs.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            raise

    def generate_slide_previews(self, output_folder="outputs/previews"):
        """Generate simple slide preview images"""
        try:
            os.makedirs(output_folder, exist_ok=True)

            for idx, slide in enumerate(self.prs.slides):
                title = slide.shapes.title.text if slide.shapes.title else f"Slide {idx+1}"
                content = ""
                
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        if shape.text_frame.text:
                            content = shape.text_frame.text.split('\n')[0]
                            break

                # Create preview image
                self._create_preview_image(idx, title, content, output_folder)

            logger.info(f"Generated {len(self.prs.slides)} previews in {output_folder}")
            
        except Exception as e:
            logger.error(f"Error generating previews: {e}")

    def _create_preview_image(self, idx, title, content, output_folder):
        """Create a single preview image"""
        try:
            # Create image with gradient background
            img = Image.new('RGB', (400, 225), color=(240, 240, 240))
            draw = ImageDraw.Draw(img)
            
            # Add gradient background
            for y in range(225):
                r = 240 - y//2
                g = 240 - y//3
                b = 240 - y//4
                draw.line([(0, y), (400, y)], fill=(r, g, b))
            
            # Try to load font
            font = None
            font_paths = [
                "assets/fonts/arial.ttf",
                "arial.ttf",
                "/System/Library/Fonts/Helvetica.ttc",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
            ]
            
            for fp in font_paths:
                try:
                    title_font = ImageFont.truetype(fp, 16)
                    content_font = ImageFont.truetype(fp, 12)
                    font = True
                    break
                except:
                    continue
            
            if not font:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
            
            # Draw slide number
            draw.rectangle([(10, 10, 50, 35)], fill=(70, 130, 180), outline=(50, 110, 160), width=2)
            draw.text((20, 15), f"{idx+1}", fill=(255, 255, 255), font=content_font, anchor="mm")
            
            # Draw title
            title_lines = self._wrap_text(title, 30)
            for i, line in enumerate(title_lines[:2]):
                draw.text((60, 15 + i*18), line, fill=(0, 51, 102), font=title_font)
            
            # Draw content preview
            if content:
                content_lines = self._wrap_text(content, 45)
                for i, line in enumerate(content_lines[:3]):
                    draw.text((15, 60 + i*16), line, fill=(51, 51, 51), font=content_font)
            
            # Add border
            draw.rectangle([(0, 0, 399, 224)], outline=(200, 200, 200), width=2)
            
            # Save
            preview_path = os.path.join(output_folder, f"slide_{idx+1:03d}.png")
            img.save(preview_path)
            
        except Exception as e:
            logger.error(f"Error creating preview image for slide {idx+1}: {e}")

    def _wrap_text(self, text: str, max_width: int) -> list:
        """Wrap text to fit within max_width characters"""
        if not text:
            return []
            
        words = text.split()
        lines = []
        current_line = []

        for word in words:
            if len(' '.join(current_line + [word])) <= max_width:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]

        if current_line:
            lines.append(' '.join(current_line))
            
        return lines