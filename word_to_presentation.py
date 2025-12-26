# word_to_presentation.py
from docx import Document
import logging

from src.document_processor import DocumentProcessor
from src.presentation_builder import PresentationBuilder

logger = logging.getLogger(__name__)


def convert_word_to_pptx(input_path, output_path):
    """Convert a Word (.docx) file to a PowerPoint (.pptx) presentation.

    Uses `DocumentProcessor` to extract structure and `PresentationBuilder`
    to create AI-enhanced, colorful slides. Falls back to a simple
    converter if anything goes wrong.
    """
    try:
        # Try the richer pipeline first
        dp = DocumentProcessor(input_path)
        headings, sections = dp.extract_document_structure()

        builder = PresentationBuilder()
        builder.add_slides_from_document(headings, sections)
        builder.save(output_path)
        print(f"Saved {output_path}")
        return

    except Exception as e:
        logger.warning(f"AI-enhanced conversion failed, falling back: {e}")

    # Fallback: basic conversion (keeps previous behavior)
    try:
        doc = Document(input_path)
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if getattr(para.style, 'name', '').startswith('Heading'):
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = text
                else:
                    # create a title textbox if layout doesn't have title
                    try:
                        slide.shapes.add_textbox(0, 0, 1, 1).text = text
                    except:
                        pass
                # clear content placeholder if exists
                try:
                    slide.placeholders[1].text = ''
                except:
                    pass
            else:
                if len(prs.slides) == 0:
                    slide = prs.slides.add_slide(slide_layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = "Slide"
                else:
                    slide = prs.slides[-1]

                try:
                    content_placeholder = slide.placeholders[1]
                    if content_placeholder.text:
                        content_placeholder.text += '\n' + text
                    else:
                        content_placeholder.text = text
                except Exception:
                    # if placeholder not available, try to append to notes
                    try:
                        slide.notes_slide.notes_text_frame.text += '\n' + text
                    except Exception:
                        pass

        prs.save(output_path)
        print(f"Saved {output_path} (fallback)")

    except Exception as e:
        logger.error(f"Failed to convert document: {e}")
        raise
